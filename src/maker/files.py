from __future__ import annotations

import mimetypes
import re
import zipfile
from dataclasses import dataclass
from pathlib import Path
from typing import Any
from uuid import uuid4

from fastapi import UploadFile

from .config import Settings

SUPPORTED_SUFFIXES = {".pdf", ".docx", ".pptx", ".txt", ".md", ".html", ".zip"}


@dataclass(slots=True)
class StoredUpload:
    original_name: str
    media_type: str
    stored_path: str
    extracted_text_path: str | None
    extraction_status: str
    metadata_json: dict[str, Any]


def sanitize_filename(filename: str) -> str:
    cleaned = re.sub(r"[^A-Za-z0-9._-]+", "-", filename.strip())
    return cleaned.strip(".-") or "upload"


def is_supported_filename(filename: str) -> bool:
    return Path(filename).suffix.lower() in SUPPORTED_SUFFIXES


async def store_upload(
    *,
    project_id: int,
    upload: UploadFile,
    settings: Settings,
) -> StoredUpload:
    original_name = upload.filename or "upload"
    safe_name = sanitize_filename(original_name)
    media_type = upload.content_type or mimetypes.guess_type(original_name)[0] or "application/octet-stream"

    project_upload_dir = settings.uploads_dir / f"project-{project_id}"
    project_upload_dir.mkdir(parents=True, exist_ok=True)

    token = uuid4().hex[:8]
    stored_path = project_upload_dir / f"{token}-{safe_name}"
    stored_path.write_bytes(await upload.read())

    extracted_text_path, extraction_status, metadata_json = extract_text_to_file(
        stored_path=stored_path,
        settings=settings,
        project_id=project_id,
        token=token,
    )

    return StoredUpload(
        original_name=original_name,
        media_type=media_type,
        stored_path=str(stored_path),
        extracted_text_path=str(extracted_text_path) if extracted_text_path else None,
        extraction_status=extraction_status,
        metadata_json=metadata_json,
    )


def load_extracted_text(path: str | None) -> str:
    if not path:
        return ""
    extracted_path = Path(path)
    if not extracted_path.exists():
        return ""
    return extracted_path.read_text(encoding="utf-8", errors="ignore")


def extract_text_to_file(
    *,
    stored_path: Path,
    settings: Settings,
    project_id: int,
    token: str,
) -> tuple[Path | None, str, dict[str, Any]]:
    suffix = stored_path.suffix.lower()
    metadata: dict[str, Any] = {"suffix": suffix}

    if suffix not in SUPPORTED_SUFFIXES:
        metadata["reason"] = "Unsupported file type."
        return None, "unsupported", metadata

    try:
        if suffix in {".txt", ".md"}:
            text = stored_path.read_text(encoding="utf-8", errors="ignore")
        elif suffix == ".html":
            text = _read_html_text(stored_path)
        elif suffix == ".pdf":
            text = _read_pdf_text(stored_path)
        elif suffix == ".docx":
            text = _read_docx_text(stored_path)
        elif suffix == ".pptx":
            text = _read_pptx_text(stored_path)
        else:
            text, metadata = _read_zip_text(stored_path, settings, project_id, token)

        extracted_dir = settings.extracted_dir / f"project-{project_id}"
        extracted_dir.mkdir(parents=True, exist_ok=True)
        extracted_path = extracted_dir / f"{token}-{stored_path.stem}.txt"
        extracted_path.write_text(text, encoding="utf-8")
        metadata["characters"] = len(text)
        status = "extracted" if text.strip() else "empty"
        return extracted_path, status, metadata
    except Exception as exc:  # pragma: no cover - defensive branch
        metadata["error"] = str(exc)
        return None, "error", metadata


def _read_html_text(path: Path) -> str:
    try:
        from bs4 import BeautifulSoup
    except ImportError as exc:  # pragma: no cover - depends on environment
        raise RuntimeError("beautifulsoup4 is not installed") from exc

    soup = BeautifulSoup(path.read_text(encoding="utf-8", errors="ignore"), "html.parser")
    return soup.get_text("\n", strip=True)


def _read_pdf_text(path: Path) -> str:
    try:
        from pypdf import PdfReader
    except ImportError as exc:  # pragma: no cover - depends on environment
        raise RuntimeError("pypdf is not installed") from exc

    reader = PdfReader(str(path))
    return "\n\n".join(page.extract_text() or "" for page in reader.pages)


def _read_docx_text(path: Path) -> str:
    try:
        from docx import Document
    except ImportError as exc:  # pragma: no cover - depends on environment
        raise RuntimeError("python-docx is not installed") from exc

    document = Document(str(path))
    paragraphs = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    return "\n".join(paragraphs)


def _read_pptx_text(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:  # pragma: no cover - depends on environment
        raise RuntimeError("python-pptx is not installed") from exc

    presentation = Presentation(str(path))
    text_fragments: list[str] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        slide_text: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and str(shape.text).strip():
                slide_text.append(str(shape.text).strip())
        if slide_text:
            text_fragments.append(f"Slide {slide_index}\n" + "\n".join(slide_text))
    return "\n\n".join(text_fragments)


def _read_zip_text(path: Path, settings: Settings, project_id: int, token: str) -> tuple[str, dict[str, Any]]:
    extract_root = settings.uploads_dir / f"project-{project_id}" / f"{token}-unzipped"
    extract_root.mkdir(parents=True, exist_ok=True)

    with zipfile.ZipFile(path) as archive:
        archive.extractall(extract_root)

    member_chunks: list[str] = []
    members: list[dict[str, Any]] = []
    for child in sorted(extract_root.rglob("*")):
        if not child.is_file():
            continue
        suffix = child.suffix.lower()
        if suffix not in SUPPORTED_SUFFIXES or suffix == ".zip":
            members.append({"path": str(child.relative_to(extract_root)), "status": "skipped"})
            continue
        try:
            if suffix in {".txt", ".md"}:
                text = child.read_text(encoding="utf-8", errors="ignore")
            elif suffix == ".html":
                text = _read_html_text(child)
            elif suffix == ".pdf":
                text = _read_pdf_text(child)
            elif suffix == ".docx":
                text = _read_docx_text(child)
            else:
                text = _read_pptx_text(child)
            relative = str(child.relative_to(extract_root))
            member_chunks.append(f"## {relative}\n{text.strip()}")
            members.append({"path": relative, "status": "extracted", "characters": len(text)})
        except Exception as exc:  # pragma: no cover - defensive branch
            members.append({"path": str(child.relative_to(extract_root)), "status": "error", "error": str(exc)})

    combined = "\n\n".join(chunk for chunk in member_chunks if chunk.strip())
    return combined, {"members": members}

